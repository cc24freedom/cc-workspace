# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案：深灰 + 商务蓝
COLOR_DARK_GRAY = RGBColor(44, 62, 80)      # #2C3E50
COLOR_BLUE = RGBColor(52, 152, 219)         # #3498DB
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(236, 240, 241)
FONT_NAME = "Microsoft YaHei"

# 图片路径
IMG_DIR = r"e:\CC\cc-tools\tools\generated_images"
IMAGES = {
    1: "gemini_1772783657.jpg",
    2: "gemini_1772783697.jpg",
    3: "gemini_1772784187.jpg",
    4: "gemini_1772783473.jpg",
    5: "gemini_1772784207.jpg",
    6: "gemini_1772783579.jpg",
    7: "gemini_1772783500.jpg",
    8: "gemini_1772784133.jpg",
    9: "gemini_1772783604.jpg",
    10: "gemini_1772783802.jpg",
    11: "gemini_1772784259.jpg",
    12: "gemini_1772783525.jpg"
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_white_background(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def add_title(slide, text, top=0.5, size=36, color=COLOR_DARK_GRAY):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_text_box(slide, text, left, top, width, height, size=18, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT, bold=False):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = text_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = align
    if bold:
        tf.paragraphs[0].font.bold = True

def add_image(slide, img_num, left, top, width):
    img_path = os.path.join(IMG_DIR, IMAGES[img_num])
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))

# 第1页：封面
def slide_1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_GRAY
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    add_title(slide, "AI Gamebot 游戏平衡性测试解决方案", 2.5, 44, COLOR_WHITE)
    add_text_box(slide, "面向 MMO/MOBA 游戏的新内容平衡验证", 2, 3.5, 6, 0.5, 24, COLOR_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, "超参数科技（深圳）有限公司\n2026年3月", 3, 6, 4, 0.8, 16, COLOR_LIGHT_GRAY, PP_ALIGN.CENTER)

# 第2页：核心价值
def slide_2_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "传统测试的三大瓶颈", 0.5, 32)

    add_text_box(slide, "样本量不足", 1, 1.8, 2.5, 0.4, 22, COLOR_DARK_GRAY, bold=True)
    add_text_box(slide, "人工测试：数十至数百场（1-2周）\nAI测试：百万场级别（1-3天）", 1, 2.5, 2.5, 0.8, 16)

    add_text_box(slide, "水平单一", 3.8, 1.8, 2.5, 0.4, 22, COLOR_DARK_GRAY, bold=True)
    add_text_box(slide, "依赖测试员水平，难以覆盖\n新手到高端全段位", 3.8, 2.5, 2.5, 0.8, 16)

    add_text_box(slide, "决策能力受限", 6.6, 1.8, 2.5, 0.4, 22, COLOR_DARK_GRAY, bold=True)
    add_text_box(slide, "脚本Bot行为固定，无法模拟\n真实玩家的适应性决策", 6.6, 2.5, 2.5, 0.8, 16)

    add_image(slide, 2, 2.25, 4.2, 5.5)

# 第3页：应用场景概览
def slide_3_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "四大核心应用场景", 0.5, 32)
    add_image(slide, 3, 2.25, 2, 5.5)

# 第4页：场景一
def slide_4_scenario1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "新英雄平衡验证", 0.5, 32)
    add_text_box(slide, "在上线前完成数值收敛", 3, 1.3, 4, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 4, 2.25, 2.2, 5.5)
    add_text_box(slide, "关键测试维度：胜率、对线强度、团战贡献、\n发育曲线、装备适配性、Counter关系", 5.5, 3, 4, 1, 16)

# 第5页：场景二
def slide_5_scenario2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "新装备/道具平衡验证", 0.5, 32)
    add_text_box(slide, "一件装备影响数十个英雄", 3, 1.3, 4, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 5, 2.25, 2.2, 5.5)
    add_text_box(slide, "输出数据：装备采用率、胜率影响、\n高风险英雄列表、装备组合预警", 1, 5.5, 8, 0.8, 16, align=PP_ALIGN.CENTER)

# 第6页：场景三
def slide_6_scenario3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "技能重做/数值调整验证", 0.5, 32)
    add_text_box(slide, "预测调整影响，避免过度削弱或强化", 2.5, 1.3, 5, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 6, 2.25, 2.2, 5.5)

# 第7页：场景四
def slide_7_scenario4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "赛季 Meta 预测与调控", 0.5, 32)
    add_text_box(slide, "提前识别\"版本答案\"", 3.5, 1.3, 3, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 7, 2.25, 2.2, 5.5)

# 第8页：技术优势一
def slide_8_advantage1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "已有成熟技术基础", 0.5, 32)
    add_text_box(slide, "基于多智能体对抗项目的完整技术积累", 2.5, 1.3, 5, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 8, 2.25, 2.2, 5.5)
    add_text_box(slide, "✓ 交付周期更短  ✓ 技术风险更低  ✓ 方案可信度更高", 1.5, 6, 7, 0.5, 18, align=PP_ALIGN.CENTER)

# 第9页：技术优势二
def slide_9_advantage2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "水平分层能力", 0.5, 32)
    add_text_box(slide, "覆盖全段位测试需求", 3.5, 1.3, 3, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 9, 2.25, 2.2, 5.5)

# 第10页：技术优势三
def slide_10_advantage3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "并行规模化，统计级样本量", 0.5, 32)
    add_text_box(slide, "百万场级别数据，充分统计意义", 2.5, 1.3, 5, 0.4, 18, COLOR_BLUE, PP_ALIGN.CENTER)
    add_image(slide, 10, 2.25, 2.2, 5.5)
    add_text_box(slide, "人工测试：数十至数百场（1-2周）\nAI Gamebot：数十万至百万场（1-3天）", 2, 5.5, 6, 0.8, 18, align=PP_ALIGN.CENTER)

# 第11页：技术优势四
def slide_11_advantage4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "技术路线已被行业验证", 0.5, 32)
    add_image(slide, 11, 2.25, 2, 5.5)
    add_text_box(slide, "EA SEED、DeepMind AlphaStar、OpenAI Five", 2, 4.8, 6, 0.5, 18, align=PP_ALIGN.CENTER)
    add_text_box(slide, "已有多智能体对抗场景的完整技术积累，可快速适配", 1.5, 6, 7, 0.6, 16, align=PP_ALIGN.CENTER)

# 第12页：延伸应用
def slide_12_extension(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)
    add_title(slide, "延伸应用与合作推进", 0.5, 32)
    add_image(slide, 12, 2.25, 2, 5.5)
    add_text_box(slide, "延伸应用：掉线玩家替补、匹配质量优化、\n新手保护机制、PvE内容支持", 1, 5, 8, 0.8, 16, align=PP_ALIGN.CENTER)

def main():
    prs = create_presentation()
    slide_1_cover(prs)
    slide_2_value(prs)
    slide_3_scenarios(prs)
    slide_4_scenario1(prs)
    slide_5_scenario2(prs)
    slide_6_scenario3(prs)
    slide_7_scenario4(prs)
    slide_8_advantage1(prs)
    slide_9_advantage2(prs)
    slide_10_advantage3(prs)
    slide_11_advantage4(prs)
    slide_12_extension(prs)

    output_path = r"e:\CC\项目\坦克世界-平衡性测试\AI_Gamebot_方案_v4.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")

if __name__ == "__main__":
    main()

