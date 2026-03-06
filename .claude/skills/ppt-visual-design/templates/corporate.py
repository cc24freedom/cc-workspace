# -*- coding: utf-8 -*-
"""商务风格 PPT 模板"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from .color_schemes import get_scheme

FONT_NAME = "Microsoft YaHei"

def corporate_cover(prs, title, subtitle, company, variant="conservative"):
    """商务风格封面页"""
    colors = get_scheme("corporate", variant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["white"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # 顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2))
    header.fill.solid()
    header.fill.fore_color.rgb = colors["primary"]
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
        sf = sub_box.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(20)
        sf.paragraphs[0].font.color.rgb = colors["accent"]
        sf.paragraphs[0].font.name = FONT_NAME
        sf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 公司信息
    if company:
        comp_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        cf = comp_box.text_frame
        cf.text = company
        cf.paragraphs[0].font.size = Pt(16)
        cf.paragraphs[0].font.color.rgb = colors["accent"]
        cf.paragraphs[0].font.name = FONT_NAME
        cf.paragraphs[0].alignment = PP_ALIGN.LEFT

def corporate_data_slide(prs, title, data_points, variant="conservative"):
    """商务风格数据页：标题+要点列表"""
    colors = get_scheme("corporate", variant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["white"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 数据要点
    y_pos = 1.5
    for point in data_points[:5]:  # 最多5个要点
        point_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.8))
        pf = point_box.text_frame
        pf.text = f"• {point}"
        pf.paragraphs[0].font.size = Pt(18)
        pf.paragraphs[0].font.color.rgb = colors["primary"]
        pf.paragraphs[0].font.name = FONT_NAME
        y_pos += 1
