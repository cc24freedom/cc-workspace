# -*- coding: utf-8 -*-
"""极简风格 PPT 模板"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from .color_schemes import get_scheme

FONT_NAME = "Microsoft YaHei"

def create_minimalist_slide(prs, colors):
    """创建极简风格空白页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["white"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide

def add_minimalist_title(slide, text, colors, top=0.5, size=32):
    """添加极简风格标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def minimalist_cover(prs, title, subtitle, footer, variant="business"):
    """极简风格封面页"""
    colors = get_scheme("minimalist", variant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["primary"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["white"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
        sf = sub_box.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(24)
        sf.paragraphs[0].font.color.rgb = colors["accent"]
        sf.paragraphs[0].font.name = FONT_NAME
        sf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 底部信息
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(4), Inches(0.8))
        ff = footer_box.text_frame
        ff.text = footer
        ff.paragraphs[0].font.size = Pt(16)
        ff.paragraphs[0].font.color.rgb = colors["light"]
        ff.paragraphs[0].font.name = FONT_NAME
        ff.paragraphs[0].alignment = PP_ALIGN.CENTER

def minimalist_image_center(prs, title, image_path, caption=None, variant="business"):
    """极简风格：标题+居中图片+说明"""
    colors = get_scheme("minimalist", variant)
    slide = create_minimalist_slide(prs, colors)
    add_minimalist_title(slide, title, colors)

    # 图片
    slide.shapes.add_picture(image_path, Inches(2.25), Inches(2), width=Inches(5.5))

    # 说明文字
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(0.8))
        cf = cap_box.text_frame
        cf.text = caption
        cf.paragraphs[0].font.size = Pt(16)
        cf.paragraphs[0].font.color.rgb = colors["primary"]
        cf.paragraphs[0].font.name = FONT_NAME
        cf.paragraphs[0].alignment = PP_ALIGN.CENTER
