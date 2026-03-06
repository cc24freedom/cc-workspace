# -*- coding: utf-8 -*-
"""创意风格 PPT 模板"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from .color_schemes import get_scheme

FONT_NAME = "Microsoft YaHei"

def creative_dark_cover(prs, title, subtitle, variant="neon"):
    """创意风格暗黑封面"""
    colors = get_scheme("creative", variant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 黑色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["primary"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # 超大标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(56)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["accent"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        sf = sub_box.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(24)
        sf.paragraphs[0].font.color.rgb = colors["white"]
        sf.paragraphs[0].font.name = FONT_NAME
        sf.paragraphs[0].alignment = PP_ALIGN.LEFT

def creative_asymmetric(prs, title, image_path, text, variant="neon"):
    """创意风格非对称布局"""
    colors = get_scheme("creative", variant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors.get("light", colors["white"])
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # 左侧大标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 右侧图片
    slide.shapes.add_picture(image_path, Inches(5), Inches(1), width=Inches(4.5))

    # 底部文字
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    txf = text_box.text_frame
    txf.text = text
    txf.paragraphs[0].font.size = Pt(20)
    txf.paragraphs[0].font.color.rgb = colors["accent"]
    txf.paragraphs[0].font.name = FONT_NAME
