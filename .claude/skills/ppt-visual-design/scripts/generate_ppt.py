#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 生成脚本
使用 python-pptx 生成技术汇报 PPT 模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os
import sys

def hex_to_rgb(hex_color):
    """将十六进制颜色转换为 RGB"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_ppt(project_name, color_scheme, pages, images_dir, output_path):
    """
    创建 PPT 文件

    Args:
        project_name: 项目名称
        color_scheme: 配色方案 {"primary": "#xxx", "secondary": "#xxx", "accent": "#xxx", "bg": "#xxx"}
        pages: 页面列表 [{"title": "xxx", "layout": "cover|content|data", "content": "xxx"}, ...]
        images_dir: 图片目录
        output_path: 输出路径
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 解析配色
    primary = RGBColor(*hex_to_rgb(color_scheme['primary']))
    secondary = RGBColor(*hex_to_rgb(color_scheme['secondary']))
    accent = RGBColor(*hex_to_rgb(color_scheme['accent']))
    bg = RGBColor(*hex_to_rgb(color_scheme['bg']))

    for page in pages:
        layout_type = page.get('layout', 'content')

        if layout_type == 'cover':
            slide = create_cover_slide(prs, page, primary, accent, images_dir)
        elif layout_type == 'content':
            slide = create_content_slide(prs, page, primary, secondary, accent)
        elif layout_type == 'data':
            slide = create_data_slide(prs, page, primary, accent)
        else:
            slide = create_content_slide(prs, page, primary, secondary, accent)

    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")

def create_cover_slide(prs, page, primary, accent, images_dir):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 背景图（如果有）
    cover_img = os.path.join(images_dir, 'cover_bg.png')
    if os.path.exists(cover_img):
        slide.shapes.add_picture(cover_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = page['title']
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, page, primary, secondary, accent):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = page['title']
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary

    # 内容区域（预留）
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = page.get('content', '内容区域')
    content_frame.paragraphs[0].font.size = Pt(20)

    return slide

def create_data_slide(prs, page, primary, accent):
    """创建数据展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = page['title']
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary

    # 核心数据（大字号）
    data_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    data_frame = data_box.text_frame
    data_frame.text = page.get('data', '核心数据')
    data_frame.paragraphs[0].font.size = Pt(60)
    data_frame.paragraphs[0].font.bold = True
    data_frame.paragraphs[0].font.color.rgb = accent
    data_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

if __name__ == '__main__':
    # 示例用法
    if len(sys.argv) < 2:
        print("用法: python generate_ppt.py <config.json>")
        sys.exit(1)

    with open(sys.argv[1], 'r', encoding='utf-8') as f:
        config = json.load(f)

    create_ppt(
        config['project_name'],
        config['color_scheme'],
        config['pages'],
        config.get('images_dir', './images'),
        config['output_path']
    )
